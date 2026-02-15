import io

from pptx import Presentation


def parse_pptx(data: bytes) -> list[dict]:
    """Parse PPTX and return list of {text, metadata} per slide."""
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if texts:
            slides.append({"text": "\n".join(texts), "metadata": {"slide": i + 1}})
    return slides
