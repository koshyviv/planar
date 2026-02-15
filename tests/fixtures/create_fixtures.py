#!/usr/bin/env python3
"""Generate minimal test fixture files for Planar tests.

Run once:  python3 tests/fixtures/create_fixtures.py
"""
import csv
import io
import os

DIR = os.path.dirname(os.path.abspath(__file__))


def create_csv():
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["Name", "Revenue", "Quarter"])
    writer.writerow(["Product A", "1200000", "Q1"])
    writer.writerow(["Product B", "850000", "Q1"])
    writer.writerow(["Product A", "1400000", "Q2"])
    with open(os.path.join(DIR, "sample.csv"), "w") as f:
        f.write(buf.getvalue())


def create_txt():
    with open(os.path.join(DIR, "sample.txt"), "w") as f:
        f.write("This is a sample text document for testing the Planar ingestion pipeline.\n")
        f.write("It contains multiple sentences to verify chunking behavior.\n")
        f.write("The quick brown fox jumps over the lazy dog.\n")


def create_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "Test Presentation"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Sample fixture for Planar tests"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.placeholders[0].text = "Slide Two"
    if len(slide2.placeholders) > 1:
        slide2.placeholders[1].text = "Some bullet points about testing."
    prs.save(os.path.join(DIR, "sample.pptx"))


def create_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("Sample Document", level=1)
    doc.add_paragraph("This is a test document for the Planar ingestion pipeline.")
    doc.add_paragraph("It tests DOCX parsing with multiple paragraphs and sections.")
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Additional content in the second section of the document.")
    doc.save(os.path.join(DIR, "sample.docx"))


def create_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Product", "Q1", "Q2", "Q3"])
    ws.append(["Widget A", 10000, 12000, 15000])
    ws.append(["Widget B", 8000, 9500, 11000])
    ws2 = wb.create_sheet("Expenses")
    ws2.append(["Category", "Amount"])
    ws2.append(["Marketing", 5000])
    ws2.append(["Engineering", 20000])
    wb.save(os.path.join(DIR, "sample.xlsx"))


def create_pdf():
    # Minimal valid PDF with text
    content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj

2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj

3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj

4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 72 720 Td (Sample PDF for Planar tests.) Tj ET
endstream
endobj

5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj

xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n

trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""
    with open(os.path.join(DIR, "sample.pdf"), "wb") as f:
        f.write(content)


if __name__ == "__main__":
    create_csv()
    create_txt()
    create_pdf()
    try:
        create_pptx()
        print("Created sample.pptx")
    except ImportError:
        print("Skipped sample.pptx (python-pptx not installed)")
    try:
        create_docx()
        print("Created sample.docx")
    except ImportError:
        print("Skipped sample.docx (python-docx not installed)")
    try:
        create_xlsx()
        print("Created sample.xlsx")
    except ImportError:
        print("Skipped sample.xlsx (openpyxl not installed)")
    print("Created sample.csv, sample.txt, sample.pdf")
