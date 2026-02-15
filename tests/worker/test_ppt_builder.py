import io

from pptx import Presentation as PptxPresentation
from app.ppt_builder import build_pptx


def test_build_pptx_basic(sample_outline):
    result = build_pptx(sample_outline)
    assert isinstance(result, bytes)
    assert len(result) > 0

    # Verify it's a valid PPTX
    prs = PptxPresentation(io.BytesIO(result))
    # Title slide + 3 content slides
    assert len(prs.slides) == 4


def test_build_pptx_empty_outline():
    outline = {"title": "Empty", "slides": []}
    result = build_pptx(outline)
    prs = PptxPresentation(io.BytesIO(result))
    assert len(prs.slides) == 1  # Just title slide


def test_build_pptx_slide_limit():
    outline = {
        "title": "Many Slides",
        "slides": [
            {"title": f"Slide {i}", "bullet_points": ["Point"], "speaker_notes": ""}
            for i in range(50)
        ],
    }
    result = build_pptx(outline)
    prs = PptxPresentation(io.BytesIO(result))
    assert len(prs.slides) <= 31  # 1 title + max 30 content
